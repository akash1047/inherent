"""Text extraction activity for converting document content to plain text.

Fetches file content directly from storage (instead of receiving bytes
via gRPC) and writes extracted text to the staging table.

Extraction dispatch (#117)
---------------------------
Which function handles a content type used to be an if/elif chain here,
duplicating the allow-list REST/MCP validation maintained independently in
``inh-public-api-svc``. Dispatch is now driven by the shared
``inh_contracts.FILE_TYPE_REGISTRY`` (the same registry REST/MCP validate
against): ``_resolve_extractor`` looks up the registry entry for a content
type and then the function wired for it in ``EXTRACTORS`` below. Two
failure modes are explicit and tested (see ``test_temporal_activities.py::
TestFileTypeRegistryDispatch``), never a silent lossy decode:

- No registry entry for the content type at all -> the document fails with
  a message naming the type and the supported set.
- A registry entry exists but its ``extractor`` key has no function in
  ``EXTRACTORS`` -- a wiring bug (a sibling format issue added a
  ``FileTypeSpec`` without its extractor) -- fails with a message that says
  so, instead of a bare ``KeyError`` crashing the Temporal worker.
"""

import datetime
import io
from collections.abc import Callable

import charset_normalizer
import structlog
from inh_contracts.file_types import all_mime_types, get_spec_for_mime
from temporalio import activity
from temporalio.exceptions import ApplicationError

from src.temporal.lineage import track_event
from src.temporal.models import ExtractTextInput, ExtractTextOutput

logger = structlog.get_logger(__name__)


@activity.defn
async def extract_text(input: ExtractTextInput) -> ExtractTextOutput:
    """Extract text from document content based on content type.

    Which formats are supported, and which function handles each, is defined
    once in ``inh_contracts.FILE_TYPE_REGISTRY`` (#117) -- see the module
    docstring above and ``docs/reference/file-types.md`` for the current
    list, rather than duplicated here where it would drift.

    The activity fetches file content from storage itself (avoiding the
    4MB gRPC limit) and writes extracted text to the staging table.

    Args:
        input: Contains storage refs, content type, filename, and workflow_run_id

    Returns:
        ExtractTextOutput with text_length (text itself is in staging)
    """
    async with track_event(
        workflow_run_id=input.workflow_run_id,
        document_id=input.document_id or "",
        workspace_id=input.workspace_id,
        event_type="text_extracted",
    ):
        return await _extract_text_inner(input)


async def _extract_text_inner(input: ExtractTextInput) -> ExtractTextOutput:
    """Inner implementation for text extraction (wrapped by lineage tracking)."""
    from src.temporal.shared_services import get_staging_service, get_storage_service

    # Fetch file content from storage
    storage_service = get_storage_service()

    if input.storage_backend == "local":
        content = storage_service.read_file(
            path=input.storage_path,
            backend="local",
            bucket=input.storage_bucket,
        )
    elif input.storage_backend == "gcs":
        content = storage_service.read_file(
            path=input.storage_path,
            backend="gcs",
            bucket=input.storage_bucket,
        )
    elif input.storage_backend == "s3":
        content = storage_service.read_file(
            path=input.storage_path,
            backend="s3",
            bucket=input.storage_bucket,
        )
    elif input.storage_backend == "azure":
        if input.storage_url:
            content = storage_service.read_file_from_url(input.storage_url)
        else:
            raise RuntimeError(f"Storage backend '{input.storage_backend}' requires storage_url")
    else:
        raise RuntimeError(f"Unknown storage backend: {input.storage_backend}")

    if content is None:
        raise RuntimeError("Failed to fetch document content from storage")

    content_type = input.content_type.lower()
    filename = input.original_filename.lower()

    logger.info(
        "Extracting text",
        content_type=content_type,
        filename=filename,
        content_size=len(content),
    )

    # Dispatch via the shared FILE_TYPE_REGISTRY (#117) -- see the module
    # docstring. `_resolve_extractor` raises a specific, actionable
    # RuntimeError for either "no registry entry" or "registry entry with no
    # wired extractor"; there is deliberately no catch-all decode-as-text
    # fallback anymore -- an unrecognized type must fail the document, never
    # silently produce garbled chunks.
    extractor = _resolve_extractor(content_type)
    text = extractor(content, input.original_filename)

    # Run data quality checks on extracted text
    from src.services.quality import DataQualityService

    quality = DataQualityService()
    quality_results = quality.check_extracted_text(text, input.original_filename)
    quality.log_results(quality_results, document_id="extract:" + input.workflow_run_id)
    if quality.has_critical_failure(quality_results):
        raise RuntimeError(
            f"Text quality check failed for {input.original_filename}: empty extraction"
        )

    # Strip NUL (0x00) bytes before staging (issue #84). Postgres text/varchar
    # columns cannot store the NUL byte at all -- the driver raises before the
    # query reaches the server -- so an unsanitized value fails the staging
    # write permanently. Some documents (e.g. PDFs pypdf decodes imperfectly)
    # produce embedded NUL bytes. We strip *after* the quality check so its
    # `no_binary_content` diagnostic still sees the raw signal, and before the
    # empty-text guard so a text made up entirely of NUL bytes is caught below.
    if "\x00" in text:
        null_count = text.count("\x00")
        text = text.replace("\x00", "")
        logger.warning(
            "Stripped NUL bytes from extracted text before staging",
            filename=input.original_filename,
            null_bytes_removed=null_count,
        )

    if not text.strip():
        raise RuntimeError(
            f"Text extraction produced empty result for {filename} "
            f"(content_type={content_type}, size={len(content)} bytes)"
        )

    logger.info(
        "Text extracted successfully",
        content_type=content_type,
        text_length=len(text),
    )

    # Write extracted text to staging
    staging = get_staging_service()
    staging.write_text(input.workflow_run_id, text)

    return ExtractTextOutput(text_length=len(text))


def _decode_text(content: bytes) -> str:
    """Decode raw bytes to text without ever silently dropping bytes (#117).

    1. Try strict UTF-8 first: the overwhelming common case for uploaded
       text/markdown/csv/html content, and the only decoding that's provably
       correct when it succeeds -- it also handles content with embedded NUL
       bytes exactly right (needed for the #84 NUL-stripping step above),
       where encoding-detection heuristics on short/ambiguous input can
       misfire (see ``test_temporal_activities.py::TestDecodeText``).
    2. If the bytes are NOT valid UTF-8, use charset-normalizer to detect the
       actual encoding (Windows-1252, UTF-16, ...) and decode with it. This
       replaces the previous ``content.decode("utf-8", errors="ignore")``,
       which silently DELETED every byte that wasn't valid UTF-8 -- data loss
       with no signal, not even a log line.
    3. If detection can't confidently identify an encoding either, fall back
       to UTF-8 with replacement characters (a visible mojibake marker)
       rather than silently vanishing the byte.
    """
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        pass

    best = charset_normalizer.from_bytes(content).best()
    if best is not None:
        return str(best)
    return content.decode("utf-8", errors="replace")


def _extract_json_text(content: bytes) -> str:
    """Parse JSON and pretty-print it as extractable text.

    ``json.loads`` accepts bytes directly (auto-detecting UTF-8/16/32 via any
    BOM present, per the JSON spec), so this needs no separate decode step.
    """
    import json

    data = json.loads(content)
    return json.dumps(data, indent=2)


# Extraction dispatch table (#117): one entry per FileTypeSpec.extractor key
# in inh_contracts.FILE_TYPE_REGISTRY. Adding a sibling format (#118 XLSX,
# #119 PPTX, ...) means adding ONE FileTypeSpec entry (services/inh-contracts)
# and ONE function + entry here -- `test_every_registry_extractor_key_is_wired`
# fails CI if the two ever disagree. Every extractor has the uniform
# ``(content: bytes, filename: str) -> str`` signature `_resolve_extractor`
# dispatches through, even where a given extractor ignores `filename` --
# lambdas adapt the helpers below that predate this table and only take
# `content`, so their own signatures/tests (and imports elsewhere) stay
# untouched. Lambdas also defer name lookup to call time, so this table can
# sit above the helper functions it references without a definition-order
# NameError at import.
EXTRACTORS: dict[str, Callable[[bytes, str], str]] = {
    "text_passthrough": lambda content, filename: _decode_text(content),
    "json_pretty": lambda content, filename: _extract_json_text(content),
    "html": lambda content, filename: _extract_html_text(content),
    "pdf": lambda content, filename: _extract_pdf_text(content),
    "docx": lambda content, filename: _extract_docx_text(content, filename),
    "xlsx": lambda content, filename: _extract_xlsx_text(content),
    "pptx": lambda content, filename: _extract_pptx_text(content),
    "image_ocr": lambda content, filename: _extract_image_text(content, filename),
}


def _resolve_extractor(content_type: str) -> Callable[[bytes, str], str]:
    """Resolve the extractor function for `content_type`, or fail loudly.

    Two distinct, explicit failure modes -- both DETERMINISTIC (the same
    `content_type` will fail identically on every retry, since neither is a
    transient dependency/network condition), so both raise a non-retryable
    ``ApplicationError`` instead of a bare ``RuntimeError``: Temporal fails
    the activity after the FIRST attempt rather than burning the workflow's
    full retry budget (multiple attempts with backoff) on a bug retrying
    cannot fix (#117 review item 13) -- cheaper and faster to reach the same
    terminal `failed` document status. Contrast with other RuntimeErrors in
    this module (storage read failures, a missing extraction library) that
    genuinely may succeed on retry and deliberately keep the default retry
    policy.

    1. No FILE_TYPE_REGISTRY entry for this content type at all -- an
       unsupported or unrecognized upload reaching extraction.
    2. A registry entry exists but its ``extractor`` key has no function
       wired into EXTRACTORS -- a wiring bug (a sibling format issue added a
       FileTypeSpec without its extractor), not a bad upload, but it must
       still fail the document with an actionable message instead of a bare
       KeyError crashing the Temporal worker.
    """
    spec = get_spec_for_mime(content_type)
    if spec is None:
        raise ApplicationError(
            f"No extractor registered for content type '{content_type}'. "
            f"Supported types: {', '.join(all_mime_types())}",
            type="UnregisteredContentType",
            non_retryable=True,
        )

    extractor = EXTRACTORS.get(spec.extractor)
    if extractor is None:
        raise ApplicationError(
            f"Registry entry '{spec.key}' ({content_type}) names extractor "
            f"'{spec.extractor}', which has no function wired in EXTRACTORS. "
            f"This is a wiring bug: add EXTRACTORS['{spec.extractor}'] in "
            f"src/temporal/activities/extract.py.",
            type="ExtractorWiringBug",
            non_retryable=True,
        )
    return extractor


def _extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF content.

    Raises on failure so Temporal retries the activity instead of
    silently producing an empty document.
    """
    try:
        import pypdf
    except ImportError:
        try:
            import PyPDF2 as pypdf  # type: ignore[no-redef]  # noqa: N813
        except ImportError:
            raise RuntimeError("PDF extraction libraries not available (pypdf or PyPDF2)")

    reader = pypdf.PdfReader(io.BytesIO(content))
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)


def _extract_docx_text(content: bytes, filename: str = "") -> str:
    """Extract text from DOCX content.

    Wrapped (review follow-up on #118/#119) so a mismatched OOXML sibling --
    see inh_contracts.file_types's docx entry comment: the shared ZIP magic
    across docx/xlsx/pptx means a mislabeled upload CAN reach this function
    with, say, genuine XLSX bytes -- fails with a clear, filename-bearing
    message instead of leaking python-docx's raw exception (observed
    verbatim: ``ValueError: file '<_io.BytesIO object at 0x7f...>' is not a
    Word file, content type is '...spreadsheetml.sheet.main+xml'``) -- a
    bare heap-address repr with no filename -- into the document's
    ``error_message`` and the dead-letter row. This is the extraction-stage
    safety net #118/#119's PR description names for the mislabeled-OOXML
    case; leaving it unwrapped while wrapping the two new extractors was a
    pattern-sweep miss caught on review.

    Deterministic given fixed `content` bytes -- retrying cannot change the
    outcome, so failures raise a non-retryable ``ApplicationError``, the same
    reasoning as the XLSX/PPTX open/cap failures below (and the existing
    ``_resolve_extractor`` "no extractor"/"wiring bug" cases).
    """
    try:
        from docx import Document
    except ImportError:
        raise ApplicationError(
            "python-docx not available for DOCX extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        )

    label = f" ({filename})" if filename else ""
    try:
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        # python-docx's own "wrong OOXML content type" ValueError embeds a
        # raw `<_io.BytesIO object at 0x...>` repr -- a heap address -- in
        # its message (observed verbatim: "file '<_io.BytesIO object at
        # 0x...>' is not a Word file, content type is '...'"). NEVER relay
        # `str(e)` verbatim here: extract just the useful "content type is
        # '...'" fragment for that specific, expected shape, and fall back
        # to the exception's TYPE NAME ONLY (never its message) for anything
        # else -- the safe default once one exception type from this same
        # library has already demonstrated it leaks an object repr.
        import re

        content_type_match = re.search(r"content type is '([^']+)'", str(e))
        detail = (
            f"wrong OOXML content type ({content_type_match.group(1)})"
            if content_type_match
            else type(e).__name__
        )
        raise ApplicationError(
            f"DOCX extraction failed{label}: could not read the document "
            f"({detail}). The file may be corrupt, truncated, "
            f"password-protected, or actually a different OOXML format "
            f"(e.g. XLSX or PPTX) despite its declared type.",
            type="DocxOpenFailed",
            non_retryable=True,
        ) from e


# Cost guards for XLSX extraction (#118 issue requirement: "cap evaluated
# cells (e.g. 500k) and emitted text length; exceeding -> document `failed`
# with actionable error, never OOM"). All three are checked INCREMENTALLY,
# INSIDE the row loop -- not after joining the whole workbook into one string
# -- so a pathological workbook fails fast with bounded peak memory instead
# of fully materializing before any guard can fire. The per-cell bound
# matters as much as the aggregate cap: `_MAX_XLSX_CELLS` counts CELLS, and a
# cell is an unbounded-length string -- a grid at 8% of the cell cap with
# 32KB strings in every cell reached 2.5GB peak RSS before the aggregate
# checks (measured on review) because nothing bounded any single value.
_MAX_XLSX_CELLS = 500_000
_MAX_XLSX_CELL_CHARS = 10_000
_MAX_XLSX_TEXT_CHARS = 5_000_000

# Cheap chunking insurance for #129 (which owns the real per-format chunker):
# a 10k-row sheet flattens to one string where only the FIRST chunk a
# downstream fixed-size chunker produces carries the sheet name and header
# row -- re-emitting both every N data rows means most chunks still carry
# that context even before #129 lands.
_XLSX_HEADER_REPEAT_ROWS = 50

# Mirrors the XLSX guards above for PPTX: a slide-count ceiling generous
# enough that a genuinely large deck (the #119 issue's illustrative "500
# slides" case) still extracts in full, a per-run/per-cell character bound
# so one pathological paragraph or table cell can't blow the budget alone,
# and a text-length cap checked incrementally per slide (not after joining
# the whole deck) -- same reasoning as XLSX above.
_MAX_PPTX_SLIDES = 5_000
_MAX_PPTX_RUN_CHARS = 10_000
_MAX_PPTX_TEXT_CHARS = 5_000_000


def _format_xlsx_cell(value: object) -> str:
    """Render one cell's value deterministically (#118 acceptance criterion:
    "Numbers and dates render deterministically; formula cells render
    computed values"), bounded so one pathological value cannot blow the
    text-length cost guard alone.

    ``openpyxl`` (opened with ``data_only=True``, see `_extract_xlsx_text`)
    already resolves formula cells to their last-computed value before this
    function ever sees them, so there is no formula-vs-literal branch here --
    every value arriving is already the value to render.

    Dates: openpyxl has no separate "date-only" Python type -- a cell
    formatted in Excel as a pure date (no time component) still comes back
    as a `datetime.datetime` with the time fixed at midnight, not a
    `datetime.date`. Rendering every `datetime.datetime` with a full
    ISO-8601 timestamp would put a noisy, meaningless "T00:00:00" suffix on
    every date column. Instead: a `datetime.datetime` at exactly midnight
    renders as a bare date; anything with a real time component keeps the
    full timestamp. Accepted tradeoff, stated plainly: a genuine timestamp
    that happens to land exactly on midnight also renders as date-only --
    openpyxl gives no other signal at this API layer to distinguish the two
    without dropping ``read_only=True`` for a second, non-read-only parse
    (the per-cell number-format lookup `.iter_rows(values_only=True)`
    intentionally forgoes), which would undo the memory-safety that buys us.
    """
    if value is None:
        return ""
    if isinstance(value, datetime.datetime):
        rendered = (
            value.date().isoformat()
            if value.time() == datetime.time(0, 0, 0, 0)
            else value.isoformat()
        )
    elif isinstance(value, datetime.date):
        rendered = value.isoformat()
    else:
        rendered = str(value)

    if len(rendered) > _MAX_XLSX_CELL_CHARS:
        # One pathological cell must not blow the text-length budget alone
        # -- truncate with a visible marker (never a silent partial value).
        original_len = len(rendered)
        rendered = rendered[:_MAX_XLSX_CELL_CHARS] + f"...[truncated, {original_len} chars]"
    return rendered


# Uncompressed-size gate for `_xlsx_merge_anchors` below. A worksheet's XML
# can be many orders of magnitude larger UNCOMPRESSED than the upload's
# on-disk (compressed) size -- verified on review: a 2.6MB uploaded .xlsx
# (pathological, from the BLOCKER 1 fix above: a 200x200 grid of 32KB-string
# cells written as OOXML inline strings, not shared strings) decompresses to
# a 1.3GB `sheet1.xml`. `zipfile.ZipFile.read()` has no size limit of its
# own -- reading that part unconditionally would decompress the whole 1.3GB
# into memory just to regex-search it for `<mergeCell>` tags, silently
# reintroducing the exact unbounded-memory failure mode BLOCKER 1 exists to
# close, through a completely different code path. Any worksheet whose
# uncompressed XML exceeds this cap skips merge-span annotation entirely
# (falls back to `{}`, i.e. "no marker" -- a presentation nicety, never
# worth this risk) rather than reading it. In practice this never matters
# for a legitimate large sheet either: a sheet whose iter_rows() reaches
# `_MAX_XLSX_TEXT_CHARS` (which a sheet of this raw size will hit almost
# immediately) fails via that cap regardless of whether this function ran.
_MAX_MERGE_SCAN_BYTES = 5_000_000


def _xlsx_merge_anchors(content: bytes, worksheet_path: str | None) -> dict[str, str]:
    """Best-effort ``{anchor_coordinate: "A1:D1"}`` map of `worksheet_path`'s
    merged cell ranges, so a merged cell's rendered value can carry a
    ``[merged A1:D1]`` marker instead of flattening to a value cell followed
    by silently blank cells with nothing distinguishing "merged" from
    "genuinely empty".

    Read directly from the sheet's raw XML via a targeted regex, NOT via
    openpyxl's `Worksheet.merged_cells` API -- `ReadOnlyWorksheet` (what
    ``read_only=True`` gives us) has no such attribute at all (verified:
    accessing it raises `AttributeError`).

    Gated by `_MAX_MERGE_SCAN_BYTES` on the part's UNCOMPRESSED size (see
    that constant's comment) before ever calling `zf.read()` -- checking
    `ZipInfo.file_size` costs nothing (it's zip central-directory metadata,
    already in hand once the archive is opened) and is what makes this safe
    to call unconditionally, once per sheet, regardless of the upload's
    contents.

    Returns ``{}`` (no annotation, extraction proceeds unaffected) if the
    part is too large to scan safely, or if anything else about this lookup
    fails -- merge-span annotation is a presentation nicety, never worth
    failing (or slowing, or ballooning the memory of) an otherwise-successful
    extraction over.
    """
    if not worksheet_path:
        return {}
    try:
        import re
        import zipfile

        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            info = zf.getinfo(worksheet_path)
            if info.file_size > _MAX_MERGE_SCAN_BYTES:
                return {}
            xml = zf.read(worksheet_path).decode("utf-8", errors="ignore")
        anchors: dict[str, str] = {}
        for match in re.finditer(r'<mergeCell\s+ref="([A-Z]+\d+):[A-Z]+\d+"', xml):
            full_ref = match.group(0)
            ref = re.search(r'ref="([^"]+)"', full_ref)
            if ref is None:
                continue
            full_range = ref.group(1)
            anchors[match.group(1)] = full_range
        return anchors
    except Exception:
        return {}


def _extract_xlsx_text(content: bytes) -> str:
    """Extract text from XLSX content with row-aware, sheet-boundary
    serialization (#118).

    Per sheet, emits ``## Sheet: <name>`` then one pipe-delimited line per
    non-empty row, cells in column order -- so an agent reading the
    flattened text can still tell which value sat in which column (row-aware
    serialization) and which sheet a row came from (sheet boundaries). This
    property holds for the raw extracted STRING; it does NOT by itself
    survive a downstream fixed-size chunker splitting that string apart
    (measured on review: a 10k-row sheet produces ~669 chunks, of which
    exactly one carries the header row and exactly one carries the "## Sheet:"
    line) -- `_XLSX_HEADER_REPEAT_ROWS` below re-emits both periodically as
    cheap insurance ahead of #129's real per-format chunker.

    ``data_only=True`` reads each formula cell's last-COMPUTED value (the
    cached result Excel/LibreOffice stores when it saves the file) rather
    than the formula source text -- exactly the "computed values only, no
    formula source" contract in the #118 issue. CAVEAT: a workbook whose
    formulas were never cached by a calculating engine (Excel, LibreOffice --
    including any workbook openpyxl itself wrote, since openpyxl does not
    evaluate formulas) reads every formula cell as `None` -- a formula-only
    row is then indistinguishable from a genuinely blank one and is skipped
    (see the blank-row-skip comment below); a sheet logged as "had rows but
    none evaluated to data" is the runtime signal for exactly this case.
    ``read_only=True`` streams rows instead of loading the whole workbook
    into memory, which is what makes the incremental caps below actually
    protective instead of cosmetic.

    Raises:
        ApplicationError (non_retryable=True): openpyxl can't open `content`
            at all (corrupt/truncated zip, password-protected/OLE2 file, a
            legitimately different binary format sharing the OOXML zip
            signature -- see inh_contracts.file_types's docx entry comment),
            the evaluated-cell cap is exceeded, or the running character
            count exceeds the text cap. Deterministic given fixed bytes --
            retrying cannot change the outcome, so these are non-retryable
            rather than burning Temporal's default retry budget on a
            guaranteed-repeat failure. Every message is clear and actionable
            -- never a bare zipfile/openpyxl exception surfacing to the
            caller, and never a silent partial result.
    """
    try:
        import openpyxl
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise ApplicationError(
            "openpyxl not available for XLSX extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        )

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    except Exception as e:
        # Covers corrupt/truncated zips (zipfile.BadZipFile), password-
        # protected files (OLE2/CFBF container -- not a zip at all), and any
        # other "openpyxl couldn't make sense of this" failure -- one clear
        # message instead of a library-specific exception type leaking out.
        raise ApplicationError(
            f"XLSX extraction failed: could not open workbook ({type(e).__name__}: {e}). "
            f"The file may be corrupt, truncated, password-protected, or not "
            f"actually an XLSX despite its declared type.",
            type="XlsxOpenFailed",
            non_retryable=True,
        ) from e

    sheet_parts: list[str] = []
    total_cells = 0
    total_chars = 0
    any_data_emitted = False
    try:
        for sheet in workbook.worksheets:
            heading = f"## Sheet: {sheet.title}"
            sheet_lines = [heading]
            total_chars += len(heading)
            merge_anchors = _xlsx_merge_anchors(content, getattr(sheet, "_worksheet_path", None))

            header_line: str | None = None
            rows_since_heading = 0
            rows_with_data = 0
            row_number = sheet.min_row or 1
            min_column = sheet.min_column or 1

            for row in sheet.iter_rows(values_only=True):
                total_cells += len(row)
                if total_cells > _MAX_XLSX_CELLS:
                    raise ApplicationError(
                        f"XLSX extraction failed: evaluated-cell cap "
                        f"({_MAX_XLSX_CELLS}) exceeded while reading sheet "
                        f"'{sheet.title}'. Split the workbook into smaller "
                        f"files and re-upload.",
                        type="XlsxCellCapExceeded",
                        non_retryable=True,
                    )
                # Skip fully blank rows (read-only mode yields a None-filled
                # row for a genuinely blank row, and -- see the CAVEAT above
                # -- for a row of uncached formula cells too, which looks
                # identical at this layer).
                if all(cell is None for cell in row):
                    row_number += 1
                    continue

                cells = []
                for col_offset, cell in enumerate(row):
                    rendered = _format_xlsx_cell(cell)
                    coord = f"{get_column_letter(min_column + col_offset)}{row_number}"
                    if coord in merge_anchors:
                        rendered = f"{rendered} [merged {merge_anchors[coord]}]"
                    cells.append(rendered)
                line = " | ".join(cells)

                if header_line is None:
                    header_line = line
                elif rows_since_heading >= _XLSX_HEADER_REPEAT_ROWS:
                    repeat_heading = f"## Sheet: {sheet.title} (continued)"
                    sheet_lines.append(repeat_heading)
                    sheet_lines.append(header_line)
                    total_chars += len(repeat_heading) + len(header_line)
                    rows_since_heading = 0

                sheet_lines.append(line)
                total_chars += len(line)
                rows_since_heading += 1
                rows_with_data += 1
                any_data_emitted = True
                row_number += 1

                # Checked INSIDE the row loop, right after the line that
                # pushed it over -- not after joining the whole workbook --
                # so peak memory is bounded to a small, fixed multiple of the
                # cap (a few oversized rows), never the full pathological
                # input's worth.
                if total_chars > _MAX_XLSX_TEXT_CHARS:
                    raise ApplicationError(
                        f"XLSX extraction failed: extracted text exceeds the "
                        f"{_MAX_XLSX_TEXT_CHARS}-character cap (hit while "
                        f"reading sheet '{sheet.title}'). Split the workbook "
                        f"into smaller files and re-upload.",
                        type="XlsxTextCapExceeded",
                        non_retryable=True,
                    )

            if rows_with_data == 0 and (sheet.max_row or 0) > 0:
                logger.warning(
                    "XLSX sheet had rows but none evaluated to visible data",
                    sheet=sheet.title,
                    max_row=sheet.max_row,
                    hint=(
                        "If this sheet contains formulas, data_only=True reads "
                        "only their cached computed value -- a workbook never "
                        "opened/saved by a calculating engine (Excel, "
                        "LibreOffice) has no cache, so formula cells read as "
                        "empty. Re-save the file in a spreadsheet application "
                        "before re-uploading."
                    ),
                )

            sheet_parts.append("\n".join(sheet_lines))
    finally:
        # read_only workbooks hold an open zip/file handle until closed --
        # always release it, success or failure.
        workbook.close()

    if not any_data_emitted:
        # No sheet in the entire workbook produced a real data row -- return
        # an honestly empty extraction (not "## Sheet: Sheet" masquerading as
        # content) so the caller's existing empty-extraction guard fails the
        # document instead of silently indexing a content-free one. Mirrors
        # PPTX's "0 slides -> ''" contract below -- XLSX's asymmetry (every
        # workbook has >=1 sheet, unlike a deck's 0-slides case) previously
        # meant an empty/uncached-formula workbook cleared that guard by
        # accident via its sheet-heading text alone.
        return ""

    return "\n\n".join(sheet_parts)


def _pptx_slide_title(slide: object) -> str | None:
    """Best-effort slide title, or None if this slide has no title
    placeholder (a valid, common case -- e.g. a section-divider or
    image-only slide)."""
    shapes = getattr(slide, "shapes", None)
    title_shape = getattr(shapes, "title", None) if shapes is not None else None
    if title_shape is None:
        return None
    text = (title_shape.text or "").strip()
    return text or None


def _pptx_slide_notes(slide: object) -> str | None:
    """Speaker notes text for `slide`, or None if it has no notes slide, or
    the notes slide has no non-whitespace text."""
    if not slide.has_notes_slide:
        return None
    notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
    return notes_text or None


def _pptx_bounded_text(text: str) -> str:
    """Bound a single paragraph/table-cell's rendered text so one
    pathological shape cannot blow the text-length budget alone -- mirrors
    `_format_xlsx_cell`'s per-cell bound, same reasoning (see the
    module-level cost-guard comment above `_MAX_XLSX_CELLS`)."""
    if len(text) > _MAX_PPTX_RUN_CHARS:
        original_len = len(text)
        return text[:_MAX_PPTX_RUN_CHARS] + f"...[truncated, {original_len} chars]"
    return text


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX content with slide-boundary sections (#119).

    Per slide, emits ``## Slide <n>: <title>`` (or ``## Slide <n>`` when the
    slide has no title placeholder), then every text-frame shape's text in
    shape order (reading order as authored), then any table shape's rows
    pipe-delimited (same row-aware convention as XLSX, #118), then speaker
    notes under a ``Notes:`` line -- so a query matching only speaker-notes
    text still lands in the same chunk as its slide's visible content once
    #129's chunker splits on these section boundaries. Embedded images are
    deliberately excluded in v1 (#119: "no OCR to manage costs" -- consistent
    with PNG's OCR being an explicit opt-in optional extra elsewhere in this
    module, not a default-on cost for every upload).

    Raises:
        ApplicationError (non_retryable=True): python-pptx can't open
            `content` at all (corrupt/truncated zip, password-protected/OLE2
            file, a different binary format sharing the OOXML zip
            signature), the slide-count cap is exceeded, or the running
            character count exceeds the text cap. Same "deterministic ->
            non-retryable, clear, actionable, never silent" contract as
            XLSX above.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ApplicationError(
            "python-pptx not available for PPTX extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        )

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as e:
        # Covers corrupt/truncated zips, password-protected (OLE2/CFBF)
        # files, and any other "python-pptx couldn't open this" failure.
        raise ApplicationError(
            f"PPTX extraction failed: could not open presentation "
            f"({type(e).__name__}: {e}). The file may be corrupt, truncated, "
            f"password-protected, or not actually a PPTX despite its "
            f"declared type.",
            type="PptxOpenFailed",
            non_retryable=True,
        ) from e

    slide_parts: list[str] = []
    total_chars = 0
    for index, slide in enumerate(presentation.slides, start=1):
        if index > _MAX_PPTX_SLIDES:
            raise ApplicationError(
                f"PPTX extraction failed: slide cap ({_MAX_PPTX_SLIDES}) "
                f"exceeded. Split the deck into smaller files and re-upload.",
                type="PptxSlideCapExceeded",
                non_retryable=True,
            )

        title = _pptx_slide_title(slide)
        heading = f"## Slide {index}: {title}" if title else f"## Slide {index}"
        slide_lines = [heading]
        total_chars += len(heading)

        for shape in slide.shapes:
            if shape.has_text_frame:
                # Title text is already in the heading above -- skip it here
                # so it isn't duplicated in the body.
                if shape == slide.shapes.title:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = _pptx_bounded_text(
                        "".join(run.text for run in paragraph.runs)
                    )
                    if paragraph_text.strip():
                        slide_lines.append(paragraph_text)
                        total_chars += len(paragraph_text)
            elif shape.has_table:
                for row in shape.table.rows:
                    line = " | ".join(_pptx_bounded_text(cell.text) for cell in row.cells)
                    slide_lines.append(line)
                    total_chars += len(line)

            # Checked INSIDE the shape loop, right after the text that pushed
            # it over -- not after joining the whole deck -- so peak memory
            # is bounded, same reasoning as XLSX's per-row check above.
            if total_chars > _MAX_PPTX_TEXT_CHARS:
                raise ApplicationError(
                    f"PPTX extraction failed: extracted text exceeds the "
                    f"{_MAX_PPTX_TEXT_CHARS}-character cap (hit on slide "
                    f"{index}). Split the deck into smaller files and "
                    f"re-upload.",
                    type="PptxTextCapExceeded",
                    non_retryable=True,
                )

        notes = _pptx_slide_notes(slide)
        if notes:
            notes = _pptx_bounded_text(notes)
            slide_lines.append("Notes:")
            slide_lines.append(notes)
            total_chars += len(notes) + len("Notes:")
            if total_chars > _MAX_PPTX_TEXT_CHARS:
                raise ApplicationError(
                    f"PPTX extraction failed: extracted text exceeds the "
                    f"{_MAX_PPTX_TEXT_CHARS}-character cap (hit on slide "
                    f"{index}'s notes). Split the deck into smaller files "
                    f"and re-upload.",
                    type="PptxTextCapExceeded",
                    non_retryable=True,
                )

        slide_parts.append("\n".join(slide_lines))

    return "\n\n".join(slide_parts)


def _extract_image_text(content: bytes, original_filename: str) -> str:
    """Extract text from a PNG image via Tesseract OCR with graceful fallback.

    OCR is an optional capability (requires the ``ocr`` extra plus the
    ``tesseract`` system binary). When OCR is unavailable -- the libraries
    are not installed, the tesseract binary is missing, or the image simply
    contains no readable text -- this returns a minimal placeholder instead
    of raising. The placeholder keeps the document flowing through the
    pipeline (0 useful chunks, but not a hard failure) so a missing OCR
    install never crashes ingestion.

    Args:
        content: Raw PNG bytes.
        original_filename: Original filename, used in the fallback placeholder.

    Returns:
        OCR-extracted text, or a placeholder string when OCR yields nothing
        or is unavailable.
    """
    placeholder = f"[image: {original_filename}, no text extracted]"

    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        logger.warning(
            "OCR libraries not available (install the 'ocr' extra: pytesseract, pillow); "
            "returning placeholder for image",
            filename=original_filename,
        )
        return placeholder

    try:
        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image)
    except pytesseract.TesseractNotFoundError:
        logger.warning(
            "Tesseract binary not found; install the 'tesseract-ocr' system package. "
            "Returning placeholder for image",
            filename=original_filename,
        )
        return placeholder
    except Exception as e:
        logger.warning(
            "OCR failed for image; returning placeholder",
            filename=original_filename,
            error=str(e),
        )
        return placeholder

    if not text.strip():
        logger.warning(
            "OCR produced no text for image; returning placeholder",
            filename=original_filename,
        )
        return placeholder

    return text


def _extract_html_text(content: bytes) -> str:
    """Extract text from HTML content.

    Falls back to a raw text decode if BeautifulSoup is not available (bs4 is
    a core, non-optional dependency of this service, so this branch is
    defense-in-depth rather than an expected runtime path). Uses
    `_decode_text` (#117) rather than `errors="ignore"` so even this fallback
    never silently drops bytes.
    """
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")
        for element in soup(["script", "style"]):
            element.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        logger.warning("beautifulsoup4 not available, falling back to raw decode")
        return _decode_text(content)
