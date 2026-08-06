"""Per-file-type extraction tests against the bundled sample documents.

Verifies that each end-to-end supported format (txt, md, csv, json, html, docx,
pdf, xlsx, pptx) extracts to non-empty, readable text via the production
extractor helpers. These run offline (no storage/staging/Temporal) by calling
the extractor helpers directly.

XLSX (#118) and PPTX (#119) were previously hard-rejected here
(``test_xlsx_extraction_is_rejected``, now replaced by ``test_extract_xlsx``
below per the #118 acceptance criterion: "flip the XLSX-rejection test... to
a positive extraction test").
"""

import io
import json
from pathlib import Path

import openpyxl
import pytest
from pptx import Presentation

from src.temporal.activities.extract import (
    _extract_docx_text,
    _extract_html_text,
    _extract_pdf_text,
    _extract_pptx_text,
    _extract_xlsx_text,
)

# tests/ -> inh-ingestion-svc -> services -> repo
_REPO_ROOT = Path(__file__).resolve().parents[3]
SAMPLE_DOCS_DIR = _REPO_ROOT / "docs" / "examples" / "sample-documents"


@pytest.fixture(autouse=True)
def cleanup_test_data():
    """Override the DB-backed root autouse fixture so these stay offline.

    These extraction tests need neither PostgreSQL nor any other live service;
    shadowing the root ``cleanup_test_data`` (which skips when PostgreSQL is
    down) lets them run unconditionally, including in local dev without Docker.
    """
    yield


def _read(filename: str) -> bytes:
    path = SAMPLE_DOCS_DIR / filename
    if not path.is_file():
        pytest.skip(f"Sample fixture missing: {path}")
    return path.read_bytes()


def test_extract_plain_text():
    text = _read("sample.txt").decode("utf-8", errors="ignore")
    assert text.strip()
    assert "Inherent" in text


def test_extract_markdown():
    text = _read("sample.md").decode("utf-8", errors="ignore")
    assert text.strip()


def test_extract_csv():
    text = _read("sample.csv").decode("utf-8", errors="ignore")
    assert text.strip()
    assert "," in text


def test_extract_json():
    data = json.loads(_read("sample.json").decode("utf-8"))
    pretty = json.dumps(data, indent=2)
    assert pretty.strip()
    assert isinstance(data, (dict, list))


def test_extract_html():
    text = _extract_html_text(_read("sample.html"))
    assert text.strip()
    # Tags must be stripped.
    assert "<html" not in text.lower()
    assert "<body" not in text.lower()


def test_extract_docx():
    text = _extract_docx_text(_read("sample.docx"))
    assert text.strip()
    assert "Inherent" in text


def test_extract_pdf():
    """Hand-built sample PDF must yield extractable text."""
    text = _extract_pdf_text(_read("sample.pdf"))
    assert text.strip(), "PDF extraction returned empty text"
    assert "Inherent" in text


def test_extract_xlsx():
    """#118: row-aware, sheet-boundary-preserving extraction against the
    bundled multi-sheet/merged-cell/empty-row fixture."""
    text = _extract_xlsx_text(_read("sample.xlsx"))
    assert text.strip()
    assert "Inherent" in text
    # Sheet boundaries are preserved as "## Sheet: <name>" headers -- an
    # agent reading the flattened text can still tell which sheet a row
    # came from.
    assert "## Sheet: Overview" in text
    assert "## Sheet: Notes" in text
    # Row-aware: cells stay pipe-delimited in column order, so "which value
    # sat in which column" survives the flatten to plain text.
    assert "Product | Region | Revenue" in text


def test_extract_pptx():
    """#119: slide-boundary sections, in-order text frames, table rows, and
    speaker notes, against the bundled fixture."""
    text = _extract_pptx_text(_read("sample.pptx"))
    assert text.strip()
    assert "Inherent" in text
    # Slide boundaries are preserved as "## Slide <n>[: <title>]" headers.
    assert "## Slide 1" in text
    # Speaker notes are appended under a "Notes:" section, not silently
    # dropped -- this is what makes a notes-only query able to retrieve the
    # right slide's chunk downstream.
    assert "Notes:" in text
    # A table shape's rows render pipe-delimited, same convention as XLSX.
    assert " | " in text


# ---------------------------------------------------------------------------
# Failure paths (#118/#119) -- corrupt input, password protection, legacy
# formats, and pathological size must all fail CLEARLY, never hang, OOM, or
# silently mis-parse. See the module docstring in inh_contracts/file_types.py
# and extract.py for why "no default lossy fallback" is the house rule.
# ---------------------------------------------------------------------------


class TestXlsxFailurePaths:
    def test_corrupt_truncated_zip_raises_runtime_error(self):
        """A ZIP-signature prefix followed by garbage (corrupt or truncated
        upload) must fail loudly -- openpyxl's zipfile.BadZipFile is caught
        and re-raised as an actionable RuntimeError, not left to propagate
        as a confusing low-level exception."""
        with pytest.raises(RuntimeError, match="Failed to open XLSX workbook"):
            _extract_xlsx_text(b"PK\x03\x04" + b"garbage, not a real zip central directory")

    def test_password_protected_bytes_raise_runtime_error(self):
        """Password-protected OOXML is wrapped in an OLE2/CFBF container
        (magic D0 CF 11 E0 A1 B1 1A E1), not a ZIP -- it is normally caught
        earlier by inh_contracts.sniff_content_type's magic-byte check
        (declared xlsx, bytes don't match the zip signature) before ever
        reaching this extractor. This test is the defense-in-depth layer:
        even called directly, the extractor itself must not crash
        ungracefully or hang -- it fails the same clear way as any other
        non-zip input."""
        ole2_encrypted_magic = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 32
        with pytest.raises(RuntimeError, match="Failed to open XLSX workbook"):
            _extract_xlsx_text(ole2_encrypted_magic)

    def test_empty_workbook_still_yields_sheet_boundary_text(self):
        """An 'empty' workbook is never truly content-less: openpyxl (and
        Excel) always ships at least one sheet, so extraction still emits
        that sheet's '## Sheet: <name>' boundary heading even with zero data
        rows. This is a DELIBERATE choice (preserve structure over emitting
        nothing) -- documented here so it isn't mistaken for a bug: contrast
        with PPTX below, where zero slides really does mean zero text."""
        workbook = openpyxl.Workbook()  # default: exactly one blank sheet
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert "## Sheet:" in text

    def test_many_sheets_extracts_all_without_hanging(self):
        """550 sheets (comfortably above a realistic real-world workbook,
        deliberately chosen to exceed the #118 issue's illustrative '500
        sheets' failure-path case) with a handful of cells each stays well
        under the cell cap and must extract every sheet, not silently drop
        or hang on any of them."""
        workbook = openpyxl.Workbook()
        workbook.remove(workbook.active)
        for i in range(550):
            sheet = workbook.create_sheet(title=f"S{i}")
            sheet["A1"] = f"row{i}"
            sheet["B1"] = i
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert text.count("## Sheet:") == 550
        assert "## Sheet: S0" in text
        assert "## Sheet: S549" in text

    def test_cell_cap_exceeded_fails_actionably_not_oom(self, monkeypatch):
        """The evaluated-cell cost guard (#118: 'cap evaluated cells (e.g.
        500k)... exceeding -> document failed with actionable error, never
        OOM') is exercised directly by lowering the cap rather than building
        a 500k-cell fixture (slow, and not the point of this test -- the
        point is the cap's ENFORCEMENT, not its specific threshold)."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_XLSX_CELLS", 3)

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["a", "b", "c", "d"])  # 4 cells > cap of 3
        buf = io.BytesIO()
        workbook.save(buf)

        with pytest.raises(RuntimeError, match="evaluated-cell cap"):
            _extract_xlsx_text(buf.getvalue())

    def test_dates_and_numbers_render_deterministically(self):
        """#118 acceptance criterion: numbers and dates render
        deterministically (not locale-dependent repr, not float noise)."""
        import datetime

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["Count", "When"])
        sheet.append([42, datetime.date(2026, 1, 15)])
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert "42" in text
        assert "2026-01-15" in text

    def test_legacy_xls_has_no_registry_entry(self):
        """Legacy .xls (binary OLE2 format, NOT the same format as .xlsx)
        must never be silently mis-parsed as XLSX. It has no
        FILE_TYPE_REGISTRY entry at all -- REST upload 400s before
        extraction is ever reached (pinned in inh-contracts'
        test_file_types.py::test_get_spec_for_extension_unknown_returns_none
        and the dispatch-layer test in test_temporal_activities.py)."""
        from inh_contracts.file_types import get_spec_for_mime

        assert get_spec_for_mime("application/vnd.ms-excel") is None


class TestPptxFailurePaths:
    def test_corrupt_truncated_zip_raises_runtime_error(self):
        with pytest.raises(RuntimeError, match="Failed to open PPTX presentation"):
            _extract_pptx_text(b"PK\x03\x04" + b"garbage, not a real zip central directory")

    def test_password_protected_bytes_raise_runtime_error(self):
        """Same reasoning as the XLSX case above: encrypted PPTX is OLE2,
        normally caught by sniff_content_type before reaching here; this
        pins the extractor's own defense-in-depth failure path."""
        ole2_encrypted_magic = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 32
        with pytest.raises(RuntimeError, match="Failed to open PPTX presentation"):
            _extract_pptx_text(ole2_encrypted_magic)

    def test_empty_deck_yields_empty_text(self):
        """Unlike XLSX (always >=1 sheet), a brand-new python-pptx
        Presentation() genuinely has ZERO slides -- there is no structural
        unit to emit a boundary heading for, so extraction returns "".
        The caller (extract_text activity) already treats an all-whitespace
        extraction as a hard failure with an actionable message -- this test
        pins the extractor's own honest "nothing to extract" contract that
        failure depends on."""
        presentation = Presentation()  # 0 slides by default
        buf = io.BytesIO()
        presentation.save(buf)

        text = _extract_pptx_text(buf.getvalue())
        assert text == ""

    def test_many_slides_extracts_all_without_hanging(self):
        """520 slides (above the issue's illustrative '500 slides' case),
        each with just a title, must all be extracted -- not silently
        truncated or hung on."""
        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        for i in range(520):
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i}"
        buf = io.BytesIO()
        presentation.save(buf)

        text = _extract_pptx_text(buf.getvalue())
        assert text.count("## Slide") == 520
        assert "Slide 0" in text
        assert "Slide 519" in text

    def test_slide_cap_exceeded_fails_actionably_not_oom(self, monkeypatch):
        """Mirrors the XLSX cell cap: a pathological deck must fail with an
        actionable message rather than run away unbounded."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_PPTX_SLIDES", 2)

        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        for i in range(4):
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i}"
        buf = io.BytesIO()
        presentation.save(buf)

        with pytest.raises(RuntimeError, match="slide cap"):
            _extract_pptx_text(buf.getvalue())

    def test_legacy_ppt_has_no_registry_entry(self):
        """Legacy .ppt (binary OLE2 format) must never be silently
        mis-parsed as PPTX -- no registry entry, same contract as .xls
        above."""
        from inh_contracts.file_types import get_spec_for_mime

        assert get_spec_for_mime("application/vnd.ms-powerpoint") is None


def test_genuine_xlsx_fed_to_docx_extractor_fails_loudly_not_silently():
    """The reachable case inh-contracts' test_file_types.py::
    test_xlsx_bytes_declared_as_docx_pass_the_byte_sniff documents: a genuine
    XLSX, declared (or defaulted, e.g. no recognized extension) as DOCX,
    passes the byte-level sniff -- the shared ZIP family magic cannot tell
    them apart. This is the layer that DOES catch it: python-docx's own
    OOXML content-type check refuses to open a package whose principal part
    is a spreadsheet, not a document, raising instead of returning mangled
    or empty text. The document fails loudly (RuntimeError -> Temporal
    activity failure -> document marked 'failed'), never silently garbled."""
    xlsx_bytes = _read("sample.xlsx")
    with pytest.raises(Exception):
        _extract_docx_text(xlsx_bytes)
